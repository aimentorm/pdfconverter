# This is a sample Python script.

# Press ⌃R to execute it or replace it with your code.
# Press Double ⇧ to search everywhere for classes, files, tool windows, actions, and settings.


def print_hi(name):
    # Use a breakpoint in the code line below to debug your script.
    print(f'Hi, {name}')  # Press ⌘F8 to toggle the breakpoint.


import os
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from pptx.util import Inches


def pdf_to_text(pdf_file):
    doc = fitz.open(pdf_file)
    text = ""
    for page in doc:
        text += page.get_text()
    return text


def pdf_to_docx(pdf_file, docx_file):
    text = pdf_to_text(pdf_file)
    doc = Document()
    doc.add_paragraph(text)
    doc.save(docx_file)
    print(f"Word document saved as {docx_file}")


def pdf_to_pptx(pdf_file, pptx_file):
    text = pdf_to_text(pdf_file)
    presentation = Presentation()
    slide_layout = presentation.slide_layouts[1]

    for line in text.splitlines():
        slide = presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = line

    presentation.save(pptx_file)
    print(f"Presentation saved as {pptx_file}")


def main():
    pdf_file = input("Enter the path to the PDF file: ")

    if not os.path.exists(pdf_file):
        print("PDF file not found!")
        return

    docx_file = input("Enter the desired output .docx file name: ")
    pptx_file = input("Enter the desired output .pptx file name: ")

    pdf_to_docx(pdf_file, docx_file)
    pdf_to_pptx(pdf_file, pptx_file)


if __name__ == "__main__":
    main()

# Press the green button in the gutter to run the script.
if __name__ == '__main__':
    print_hi('PyCharm')

# See PyCharm help at https://www.jetbrains.com/help/pycharm/
